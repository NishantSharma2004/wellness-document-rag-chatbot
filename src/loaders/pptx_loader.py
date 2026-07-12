from pathlib import Path
from typing import List, Dict, Any
from pptx import Presentation
from src.loaders.base import BaseLoader
from src.utils.exceptions import LoaderException

class PptxLoader(BaseLoader):
    def load(self, filepath: Path) -> List[Dict[str, Any]]:
        self.validate_file(filepath)
        
        pages = []
        try:
            prs = Presentation(filepath)
            for idx, slide in enumerate(prs.slides):
                slide_num = idx + 1
                slide_text = []
                
                # Extract text from shape
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                        
                pages.append({
                    "text": "\n".join(slide_text),
                    "page_number": slide_num,  # Page maps to slide number for PPTX
                    "slide_number": slide_num,
                    "section": None
                })
        except Exception as e:
            raise LoaderException(f"Failed to load PPTX file {filepath}: {str(e)}") from e
            
        return pages
